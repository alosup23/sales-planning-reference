#!/usr/bin/env python3
from __future__ import annotations

import re
import sys
import textwrap
from datetime import datetime, timezone
from pathlib import Path

DOC_DEPS = Path("/tmp/sales-planning-docdeps")
if DOC_DEPS.exists():
    sys.path.insert(0, str(DOC_DEPS))

from docx import Document  # type: ignore
from docx.enum.section import WD_SECTION_START  # type: ignore
from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
from docx.oxml.ns import qn  # type: ignore
from docx.shared import Inches, Pt  # type: ignore
from PIL import Image, ImageDraw, ImageFont  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.text import PP_ALIGN  # type: ignore
from pptx.util import Inches as PptInches, Pt as PptPt  # type: ignore

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
DELIVERABLES = DOCS / "deliverables"
ASSETS = DELIVERABLES / "assets"

USER_GUIDE_MD = DOCS / "user-guide.md"
TRAINING_MD = DOCS / "training-process-overview.md"
FILE_FORMATS_MD = DOCS / "master-data-file-formats.md"
LIMITATIONS_MD = DOCS / "current-limitations-and-recommendations.md"
AWS_CONFIG_MD = DOCS / "current-uat-aws-configuration.md"
API_FLOWS_MD = DOCS / "api-endpoints-and-transaction-flows.md"

USER_GUIDE_DOCX = DELIVERABLES / "sales-planning-user-guide-phase1-uat.docx"
AWS_STACK_PPTX = DELIVERABLES / "sales-planning-aws-stack-and-transaction-flow-phase1-uat.pptx"
TRAINING_PPTX = DELIVERABLES / "sales-budget-and-planning-training-phase1-uat.pptx"
README_TXT = DELIVERABLES / "README.txt"

AWS_STACK_PNG = ASSETS / "aws-uat-stack-diagram.png"
TX_FLOW_PNG = ASSETS / "aws-transaction-flow-diagram.png"
PLAN_FLOW_PNG = ASSETS / "sales-planning-process-flow.png"


def ensure_dirs() -> None:
    DELIVERABLES.mkdir(parents=True, exist_ok=True)
    ASSETS.mkdir(parents=True, exist_ok=True)


def load_font(size: int, bold: bool = False):
    candidates = [
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf" if bold else "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/System/Library/Fonts/Supplemental/Helvetica.ttc",
        "/Library/Fonts/Arial.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size)
    return ImageFont.load_default()


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font, max_width: int) -> list[str]:
    words = text.split()
    if not words:
        return [""]
    lines: list[str] = []
    current = words[0]
    for word in words[1:]:
        trial = f"{current} {word}"
        if draw.textbbox((0, 0), trial, font=font)[2] <= max_width:
            current = trial
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def draw_box(draw: ImageDraw.ImageDraw, xy, fill, outline, title: str, body: str, title_font, body_font) -> None:
    draw.rounded_rectangle(xy, radius=24, fill=fill, outline=outline, width=3)
    x1, y1, x2, y2 = xy
    pad = 18
    draw.text((x1 + pad, y1 + pad), title, font=title_font, fill=(28, 35, 48))
    body_y = y1 + pad + 42
    max_width = (x2 - x1) - pad * 2
    for line in wrap_text(draw, body, body_font, max_width):
        draw.text((x1 + pad, body_y), line, font=body_font, fill=(63, 74, 88))
        body_y += 26


def draw_arrow(draw: ImageDraw.ImageDraw, start, end, fill=(47, 84, 150), width=6) -> None:
    draw.line([start, end], fill=fill, width=width)
    ex, ey = end
    sx, sy = start
    if abs(ex - sx) >= abs(ey - sy):
        direction = 1 if ex > sx else -1
        points = [(ex, ey), (ex - 18 * direction, ey - 10), (ex - 18 * direction, ey + 10)]
    else:
        direction = 1 if ey > sy else -1
        points = [(ex, ey), (ex - 10, ey - 18 * direction), (ex + 10, ey - 18 * direction)]
    draw.polygon(points, fill=fill)


def make_aws_stack_diagram(path: Path) -> None:
    image = Image.new("RGB", (1800, 1000), (247, 250, 252))
    draw = ImageDraw.Draw(image)
    title_font = load_font(34, bold=True)
    box_title_font = load_font(26, bold=True)
    body_font = load_font(20)
    small_font = load_font(18)

    draw.text((60, 40), "Sales Budget & Planning UAT AWS Stack", font=title_font, fill=(28, 35, 48))
    draw.text((60, 88), "Live Phase 1 deployment after ECS, private RDS, and WAF cutover", font=small_font, fill=(78, 92, 110))

    browser = (80, 250, 340, 390)
    cloudfront = (420, 180, 760, 330)
    s3 = (420, 390, 760, 520)
    alb = (860, 180, 1200, 330)
    ecs = (1300, 180, 1640, 330)
    rds = (1300, 420, 1640, 570)
    secrets = (860, 420, 1200, 540)
    entra = (860, 620, 1200, 760)
    logs = (1300, 620, 1640, 760)

    draw_box(draw, browser, (255, 255, 255), (125, 142, 161), "User Browser", "Responsive planning web app and Microsoft sign-in.", box_title_font, body_font)
    draw_box(draw, cloudfront, (232, 242, 255), (95, 145, 220), "CloudFront + WAF", "Public HTTPS entry point, static app delivery, rate limit and managed rule protection.", box_title_font, body_font)
    draw_box(draw, s3, (236, 252, 244), (78, 166, 114), "S3 Web Bucket", "Published web bundle and static assets.", box_title_font, body_font)
    draw_box(draw, alb, (250, 240, 230), (212, 153, 89), "Origin-Protected ALB", "Only accepts CloudFront traffic with the origin verification header.", box_title_font, body_font)
    draw_box(draw, ecs, (243, 232, 255), (147, 112, 219), "ECS Fargate API", "Interactive .NET 8 planning API with Entra auth, planning commands, and master-data maintenance.", box_title_font, body_font)
    draw_box(draw, rds, (255, 244, 230), (216, 146, 75), "RDS PostgreSQL", "Authoritative Phase 1 persistence in true private subnets.", box_title_font, body_font)
    draw_box(draw, secrets, (255, 248, 220), (210, 177, 64), "Secrets Manager", "Database credentials resolved at runtime by the API.", box_title_font, body_font)
    draw_box(draw, entra, (234, 246, 255), (98, 164, 214), "Microsoft Entra ID", "Work-account sign-in and backend API audience/scope authorization.", box_title_font, body_font)
    draw_box(draw, logs, (241, 245, 249), (148, 163, 184), "CloudWatch Logs", "API task logs, health checks, and deployment diagnostics.", box_title_font, body_font)

    draw_arrow(draw, (340, 320), (420, 255))
    draw_arrow(draw, (590, 330), (590, 390))
    draw_arrow(draw, (760, 255), (860, 255))
    draw_arrow(draw, (1200, 255), (1300, 255))
    draw_arrow(draw, (1470, 330), (1470, 420))
    draw_arrow(draw, (1200, 460), (1300, 460))
    draw_arrow(draw, (1030, 620), (1300, 330))
    draw_arrow(draw, (1470, 330), (1470, 620))

    image.save(path)


def make_transaction_flow_diagram(path: Path) -> None:
    image = Image.new("RGB", (1800, 1000), (250, 251, 253))
    draw = ImageDraw.Draw(image)
    title_font = load_font(34, bold=True)
    box_title_font = load_font(24, bold=True)
    body_font = load_font(19)

    draw.text((60, 40), "Interactive Planning Transaction Flow", font=title_font, fill=(28, 35, 48))

    steps = [
        ("1. Sign In", "User opens the app, signs in with Microsoft 365, and receives a valid Entra token."),
        ("2. Load Scope", "Browser loads planning store scopes and the current grid slice through CloudFront."),
        ("3. Expand Branch", "The grid requests additional branch rows only when the user expands a hierarchy node."),
        ("4. Edit Or Splash", "The API validates locks and scope, applies the change, recalculates, and records the journal."),
        ("5. Return Patch", "The server returns only the changed cells and state needed for the active slice."),
        ("6. Save / Audit", "The user saves, exports, or reviews audit history. Undo and redo replay the same journal."),
    ]

    start_x = 80
    box_w = 250
    gap = 40
    y = 290
    box_h = 220
    for index, (title, body) in enumerate(steps):
        x = start_x + index * (box_w + gap)
        draw_box(draw, (x, y, x + box_w, y + box_h), (255, 255, 255), (109, 125, 145), title, body, box_title_font, body_font)
        if index < len(steps) - 1:
            draw_arrow(draw, (x + box_w, y + box_h // 2), (x + box_w + gap, y + box_h // 2))

    draw.text((80, 590), "Rules enforced throughout:", font=box_title_font, fill=(28, 35, 48))
    rules = [
        "Parent totals must equal child totals after every committed action.",
        "Locked descendants are excluded from top-down splash allocation.",
        "Undo / redo replays the application command journal up to 30 actions.",
        "All views read from the same canonical planning facts.",
    ]
    bullet_y = 640
    bullet_font = load_font(22)
    for rule in rules:
        draw.text((100, bullet_y), f"• {rule}", font=bullet_font, fill=(63, 74, 88))
        bullet_y += 44

    image.save(path)


def make_planning_process_diagram(path: Path) -> None:
    image = Image.new("RGB", (1800, 1000), (249, 250, 251))
    draw = ImageDraw.Draw(image)
    title_font = load_font(34, bold=True)
    box_title_font = load_font(24, bold=True)
    body_font = load_font(19)

    draw.text((60, 40), "Best-Practice Sales Budget & Planning Flow", font=title_font, fill=(28, 35, 48))
    rows = [
        [
            ("Prepare Hierarchy", "Validate Department, Class, Subclass, and option values."),
            ("Load Master Data", "Import Store, Product, Inventory, Pricing, Seasonality, and Vendor files."),
            ("Confirm Year", "Generate the next year if needed and verify period structure."),
        ],
        [
            ("Bottom-Up Plan", "Use Planning - by Store for leaf edits and detailed build."),
            ("Top-Down Align", "Use Planning - by Department for cross-store and aggregate review."),
            ("Lock Approved Values", "Protect approved branches before wider redistribution."),
        ],
        [
            ("Undo / Redo Where Needed", "Use the 30-action journal to reverse recent planning changes safely."),
            ("Export And Reconcile", "Export workbooks, review audit, and confirm totals."),
            ("Prepare For Phase 2", "Keep Inventory, Pricing, Seasonality, and Vendor context current."),
        ],
    ]

    origin_x = 70
    origin_y = 160
    box_w = 500
    box_h = 180
    gap_x = 70
    gap_y = 70
    colors = [
        ((232, 242, 255), (95, 145, 220)),
        ((236, 252, 244), (78, 166, 114)),
        ((255, 244, 230), (216, 146, 75)),
    ]

    for row_index, row in enumerate(rows):
        for col_index, (title, body) in enumerate(row):
            x1 = origin_x + col_index * (box_w + gap_x)
            y1 = origin_y + row_index * (box_h + gap_y)
            x2 = x1 + box_w
            y2 = y1 + box_h
            fill, outline = colors[row_index]
            draw_box(draw, (x1, y1, x2, y2), fill, outline, title, body, box_title_font, body_font)
            if col_index < len(row) - 1:
                draw_arrow(draw, (x2, y1 + box_h // 2), (x2 + gap_x, y1 + box_h // 2), fill=(90, 103, 125))
        if row_index < len(rows) - 1:
            mid_x = origin_x + box_w + gap_x // 2
            y_start = origin_y + row_index * (box_h + gap_y) + box_h
            y_end = y_start + gap_y
            draw_arrow(draw, (mid_x, y_start), (mid_x, y_end), fill=(90, 103, 125))

    image.save(path)


def parse_markdown(path: Path):
    blocks = []
    lines = path.read_text(encoding="utf-8").splitlines()
    paragraph: list[str] = []
    bullets: list[str] = []
    numbers: list[str] = []
    in_code = False

    def flush_paragraph():
        nonlocal paragraph
        if paragraph:
            blocks.append(("paragraph", " ".join(part.strip() for part in paragraph).strip()))
            paragraph = []

    def flush_bullets():
        nonlocal bullets
        if bullets:
            blocks.append(("bullets", bullets))
            bullets = []

    def flush_numbers():
        nonlocal numbers
        if numbers:
            blocks.append(("numbers", numbers))
            numbers = []

    for line in lines:
        if line.startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)$", line)
        bullet = re.match(r"^\-\s+(.*)$", line)
        number = re.match(r"^\d+\.\s+(.*)$", line)

        if heading:
            flush_paragraph()
            flush_bullets()
            flush_numbers()
            blocks.append(("heading", len(heading.group(1)), heading.group(2).strip()))
        elif bullet:
            flush_paragraph()
            flush_numbers()
            bullets.append(bullet.group(1).strip())
        elif number:
            flush_paragraph()
            flush_bullets()
            numbers.append(number.group(1).strip())
        elif not line.strip():
            flush_paragraph()
            flush_bullets()
            flush_numbers()
        else:
            paragraph.append(line)

    flush_paragraph()
    flush_bullets()
    flush_numbers()
    return blocks


def add_markdown_to_doc(document: Document, path: Path) -> None:
    for block in parse_markdown(path):
        kind = block[0]
        if kind == "heading":
            _, level, text = block
            document.add_heading(text, level=min(level, 4))
        elif kind == "paragraph":
            document.add_paragraph(block[1])
        elif kind == "bullets":
            for item in block[1]:
                document.add_paragraph(item, style="List Bullet")
        elif kind == "numbers":
            for item in block[1]:
                document.add_paragraph(item, style="List Number")


def build_user_guide_docx() -> None:
    document = Document()
    document.core_properties.title = "Sales Budget & Planning User Guide"
    document.core_properties.subject = "Phase 1 UAT user guidance"
    document.core_properties.author = "OpenAI Codex"
    document.core_properties.created = datetime.now(timezone.utc)

    section = document.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("Sales Budget & Planning\nPhase 1 UAT User Guide")
    run.bold = True
    run.font.size = Pt(22)

    subtitle = document.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle.add_run("Comprehensive guide covering navigation, planning flow, import/export, and bottom-up / top-down rules")
    subtitle_run.italic = True
    subtitle_run.font.size = Pt(11)

    document.add_paragraph("")
    add_markdown_to_doc(document, USER_GUIDE_MD)

    document.add_section(WD_SECTION_START.NEW_PAGE)
    document.add_heading("Appendix A: Workbook Format Summary", level=1)
    add_markdown_to_doc(document, FILE_FORMATS_MD)

    document.add_section(WD_SECTION_START.NEW_PAGE)
    document.add_heading("Appendix B: Current UAT Limitations", level=1)
    add_markdown_to_doc(document, LIMITATIONS_MD)

    for paragraph in document.paragraphs:
        for run in paragraph.runs:
            run.font.name = "Arial"
            run._element.rPr.rFonts.set(qn("w:eastAsia"), "Arial")

    document.save(USER_GUIDE_DOCX)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: list[str] | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    text_frame = slide.shapes.placeholders[1].text_frame
    text_frame.clear()
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = PptPt(22)
    if notes:
        notes_text = "\n".join(notes)
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes_text


def add_image_slide(prs: Presentation, title: str, image_path: Path, subtitle: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if subtitle:
        textbox = slide.shapes.add_textbox(PptInches(0.7), PptInches(0.8), PptInches(12), PptInches(0.4))
        p = textbox.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.size = PptPt(16)
        p.font.color.rgb = RGBColor(96, 108, 122)
    slide.shapes.add_picture(str(image_path), PptInches(0.6), PptInches(1.35), width=PptInches(12.1))


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), PptInches(0.5), PptInches(1.3), PptInches(12.2), PptInches(5.6))
    table = table_shape.table
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            table.cell(row_index, col_index).text = value
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = PptPt(16)


def build_aws_presentation() -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "Sales Planning UAT AWS Stack And Transaction Flow",
        "Phase 1 deployed runtime, endpoints, security decisions, and operational considerations",
    )
    add_bullet_slide(
        prs,
        "Current Live UAT Stack",
        [
            "CloudFront + WAF for the public web entry point",
            "Static React build on S3",
            ".NET 8 interactive API on ECS Fargate",
            "RDS PostgreSQL in true private subnets",
            "Microsoft Entra sign-in and backend API authorization",
            "Origin-protected ALB behind CloudFront",
        ],
    )
    add_image_slide(
        prs,
        "Graphical AWS Stack",
        AWS_STACK_PNG,
        "Live Phase 1 UAT runtime after the ECS, private RDS, and WAF cutover.",
    )
    add_image_slide(
        prs,
        "Transaction Flow",
        TX_FLOW_PNG,
        "Startup, lazy branch loading, command processing, patch return, save, and audit flow.",
    )
    add_table_slide(
        prs,
        "Key Endpoints",
        ["Area", "Endpoint / Name", "Purpose"],
        [
            ["Public App", "https://d22xc0mfhkv9bk.cloudfront.net", "User access entry point"],
            ["Health", "https://d22xc0mfhkv9bk.cloudfront.net/health", "Operational smoke check"],
            ["Public API", "/api/v1", "Protected business API path through CloudFront"],
            ["ECS Service", "sales-planning-demo-api", "Interactive API workload"],
            ["ECS Cluster", "sales-planning-demo-cluster", "Fargate cluster"],
            ["Active DB", "sales-planning-demo-pg-private", "Private PostgreSQL instance"],
        ],
    )
    add_bullet_slide(
        prs,
        "Security And Design Decisions",
        [
            "WAF attached at CloudFront with managed rule groups and rate limiting",
            "ALB accepts traffic only from CloudFront and only when the secret origin header is present",
            "Database is not publicly accessible and runs in the dedicated private DB subnet group",
            "ECS uses Secrets Manager for DB credentials",
            "CloudFront-to-ALB HTTPS is deferred until Route 53 and ACM are available",
        ],
    )
    add_bullet_slide(
        prs,
        "Current Limitations And Recommendations",
        [
            "Grid still uses the client-side row model rather than AG Grid SSRM",
            "Workbook import and export are still synchronous",
            "CloudFront-to-ALB origin traffic is still HTTP",
            "Rollback DB should be deleted after UAT acceptance",
            "Next priority: delta recalculation, SSRM, async jobs, and ALB HTTPS",
        ],
    )
    prs.save(AWS_STACK_PPTX)


def build_training_presentation() -> None:
    prs = Presentation()
    add_title_slide(
        prs,
        "Sales Budget And Planning Process Overview",
        "Phase 1 UAT training deck for planners, merchandisers, and business reviewers",
    )
    add_bullet_slide(
        prs,
        "Training Objectives",
        [
            "Understand the planning workspaces and when to use each one",
            "Execute the end-to-end budget and planning cycle in the recommended sequence",
            "Use bottom-up edits, top-down splash, locks, and undo/redo correctly",
            "Import and export the supported workbook formats safely",
            "Understand the core planning rules and current UAT limitations",
        ],
    )
    add_image_slide(
        prs,
        "Best-Practice Planning Flow",
        PLAN_FLOW_PNG,
        "Recommended end-to-end sequence from hierarchy setup through export and reconciliation.",
    )
    add_bullet_slide(
        prs,
        "Application Workspaces",
        [
            "Planning - by Store for detailed bottom-up planning",
            "Planning - by Department for top-down and cross-store review",
            "Hierarchy Maintenance for Department / Class / Subclass control",
            "Store, Product, Inventory, Pricing, Seasonality, and Vendor maintenance workspaces",
            "Compact menu-based toolbar preserves maximum grid space",
        ],
    )
    add_bullet_slide(
        prs,
        "Bottom-Up And Top-Down Rules",
        [
            "Leaf edits should be performed at the lowest relevant level",
            "Parent totals must always equal the sum of children after recalculation",
            "Splash reallocates aggregate targets to unlocked descendants only",
            "Same-year recalculation applies throughout Phase 1",
            "Undo and redo preserve the planning invariants for up to 30 actions",
        ],
    )
    add_bullet_slide(
        prs,
        "Supported Workbook Imports And Exports",
        [
            "Branch Profile.xlsx",
            "Product Profile.xlsx",
            "Inventory Profile.xlsx",
            "Pricing Policy.xlsx",
            "Seasonality & Events.xlsx",
            "Vendor Supply Profile.xlsx",
            "Planning workbook import and export",
        ],
    )
    add_bullet_slide(
        prs,
        "Grid And UI Guidance",
        [
            "Use the workspace selector to switch functions quickly",
            "Keep hierarchies collapsed and expand only the branches being reviewed",
            "Use store view first for detail and department view second for alignment",
            "Use locks before major top-down changes",
            "Export and review after each major planning cycle",
        ],
    )
    add_bullet_slide(
        prs,
        "Phase 2-Ready Data Foundations",
        [
            "Pricing policy, inventory, seasonality, and vendor data are already in Phase 1 scope",
            "These data foundations will support future price, markdown, and forecasting recommendations",
            "Phase 2 will reuse the current transactional planning model rather than replace it",
        ],
    )
    add_bullet_slide(
        prs,
        "Current UAT Considerations",
        [
            "Imports and exports are still synchronous in Phase 1",
            "Some performance hardening remains for production",
            "CloudFront-to-ALB HTTPS is pending Route 53 and ACM",
            "Always use the exported templates when preparing workbook imports",
        ],
    )
    prs.save(TRAINING_PPTX)


def write_deliverables_readme() -> None:
    lines = [
        "Sales Planning Phase 1 UAT Deliverables",
        "",
        f"Generated on: {datetime.now().astimezone().strftime('%Y-%m-%d %H:%M:%S %Z')}",
        "",
        "Files:",
        f"- {USER_GUIDE_DOCX.name}",
        f"- {AWS_STACK_PPTX.name}",
        f"- {TRAINING_PPTX.name}",
        "",
        "Supporting source documents:",
        "- docs/user-guide.md",
        "- docs/training-process-overview.md",
        "- docs/current-uat-aws-configuration.md",
        "- docs/api-endpoints-and-transaction-flows.md",
        "- docs/current-limitations-and-recommendations.md",
        "- docs/master-data-file-formats.md",
        "",
        "Supporting generated assets:",
        f"- assets/{AWS_STACK_PNG.name}",
        f"- assets/{TX_FLOW_PNG.name}",
        f"- assets/{PLAN_FLOW_PNG.name}",
    ]
    README_TXT.write_text("\n".join(lines) + "\n", encoding="utf-8")


def main() -> None:
    ensure_dirs()
    make_aws_stack_diagram(AWS_STACK_PNG)
    make_transaction_flow_diagram(TX_FLOW_PNG)
    make_planning_process_diagram(PLAN_FLOW_PNG)
    build_user_guide_docx()
    build_aws_presentation()
    build_training_presentation()
    write_deliverables_readme()
    print(f"Generated {USER_GUIDE_DOCX}")
    print(f"Generated {AWS_STACK_PPTX}")
    print(f"Generated {TRAINING_PPTX}")


if __name__ == "__main__":
    main()
